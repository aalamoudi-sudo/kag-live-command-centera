#!/usr/bin/env python3
"""
fill_template.py — يستقبل JSON من stdin، يملأ قالب PPTX، يخرج bytes إلى stdout
"""
import sys, json, io, os
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy

def set_shape_text(shape, new_text):
    """استبدال نص shape مع الحفاظ على تنسيق أول run"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # احفظ تنسيق أول run
    first_run = None
    for para in tf.paragraphs:
        for run in para.runs:
            first_run = run
            break
        if first_run:
            break
    
    # امسح كل paragraphs عدا الأول
    from pptx.oxml.ns import qn
    from lxml import etree
    txBody = tf._txBody
    paras = txBody.findall(qn('a:p'))
    
    # إذا نص متعدد الأسطر، أنشئ paragraph لكل سطر
    lines = str(new_text).split('\n')
    
    # امسح كل paragraphs
    for p in paras:
        txBody.remove(p)
    
    for line_idx, line in enumerate(lines):
        # أنشئ paragraph جديد
        p_elem = etree.SubElement(txBody, qn('a:p'))
        r_elem = etree.SubElement(p_elem, qn('a:r'))
        
        # انسخ تنسيق الـ run الأصلي إذا وُجد
        if first_run is not None:
            rPr = first_run._r.find(qn('a:rPr'))
            if rPr is not None:
                r_elem.insert(0, deepcopy(rPr))
        
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = line

def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None

def status_label(s):
    if not s: return '—'
    if s in ['مكتملة','مكتمل','معتمدة','معتمد','ضمن المسار','Completed','Cleared']: return 'أخضر ✓'
    if s in ['قيد التنفيذ','تحت المتابعة','In Progress','Watch']: return 'أصفر'
    if s in ['معرضة للخطر','معرض للخطر','متأخرة','At Risk','متأخر']: return 'أحمر ✗'
    return s

def format_date(d):
    if not d: return ''
    try:
        return datetime.strptime(d[:10], '%Y-%m-%d').strftime('%Y/%m/%d')
    except:
        return d

def set_slide_shape(slide, name, text):
    shape = find_shape(slide, name)
    if shape:
        set_shape_text(shape, str(text) if text is not None else '—')

def today_ar():
    d = datetime.now()
    months_ar = ['يناير','فبراير','مارس','أبريل','مايو','يونيو',
                 'يوليو','أغسطس','سبتمبر','أكتوبر','نوفمبر','ديسمبر']
    days_ar = ['الاثنين','الثلاثاء','الأربعاء','الخميس','الجمعة','السبت','الأحد']
    return f"{days_ar[d.weekday()]} {d.day} {months_ar[d.month-1]} {d.year}"

def week_num():
    return datetime.now().isocalendar()[1]

def get_track(tracks, tid):
    for t in tracks:
        if t.get('track') == tid or t.get('id') == tid:
            return t
    return {}

# ============================================================
# تعبئة التقرير الشامل
# ============================================================
def fill_comprehensive(prs, state):
    tracks = state.get('tracks', [])
    items  = state.get('items', [])
    
    done   = [i for i in items if i.get('status') in ['مكتملة','مكتمل','معتمدة','معتمد','Completed','Cleared']]
    active = [i for i in items if i.get('status') in ['قيد التنفيذ','In Progress']]
    risks  = [i for i in items if i.get('type') in ['risks','مخاطرة','مخاطر'] and i.get('status') != 'مغلقة']
    
    overall = round(sum(t.get('progress',0) for t in tracks)/len(tracks)) if tracks else 0
    worst = 'أحمر' if any(t.get('status') in ['معرض للخطر','معرضة للخطر','At Risk'] for t in tracks) \
            else 'أصفر' if any(t.get('status') == 'تحت المتابعة' for t in tracks) else 'أخضر'

    # Slide 1 — الغلاف
    s1 = prs.slides[0]
    set_slide_shape(s1, 'Text 4', f"الأسبوع {week_num()} — {datetime.now().year} | الفترة: {today_ar()}")

    # Slide 2 — الملخص التنفيذي
    s2 = prs.slides[1]
    set_slide_shape(s2, 'Text 17', f"الحالة: {worst}\nنسبة إنجاز اليوم: {overall}٪")
    
    done_titles = [i.get('title','') for i in done[:3]]
    while len(done_titles) < 3: done_titles.append('—')
    active_titles = [i.get('title','') for i in active[:3]]
    while len(active_titles) < 3: active_titles.append('—')
    
    set_slide_shape(s2, 'Text 21', '\n'.join(done_titles))
    set_slide_shape(s2, 'Text 25', '\n'.join(active_titles))

    # جدول القضايا — 3 صفوف (Text 42/44/46/48/50, Text 62/64/66/68/70, Text 72/...)
    issue_shapes = [
        ('Shape 51','Text 44','Text 46','Text 48','Text 50'),
        ('Text 62','Text 64','Text 66','Text 68','Text 70'),
        ('Text 72','Text 54','Text 56','Text 58','Text 60'),
    ]
    for idx, (sq,wq,mq,aq,athq) in enumerate(issue_shapes):
        if idx < len(risks):
            r = risks[idx]
            set_slide_shape(s2, sq, r.get('title','—'))
            set_slide_shape(s2, wq, format_date(r.get('due','')))
            set_slide_shape(s2, mq, r.get('owner','—'))
            set_slide_shape(s2, aq, 'متابعة عاجلة')
            set_slide_shape(s2, athq, 'تأثير على الجدول الزمني')

    # Slide 3 — حالة المسارات
    s3 = prs.slides[2]
    track_rows = [
        # (اسم المسار shape, إنجاز أمس, اليوم, الغد, الحالة, الدعم)
        ('Text 34','Text 32','Text 30','Text 28','Text 26','Text 24'),  # مسار أ
        ('Text 46','Text 44','Text 42','Text 40','Text 38','Text 36'),  # مسار ب
        ('Text 58','Text 56','Text 54','Text 52','Text 50','Text 48'),  # مسار ج
        ('Text 70','Text 68','Text 66','Text 64','Text 62','Text 60'),  # مسار د
    ]
    track_ids = ['أ','ب','ج','د']
    for tidx, (nm,s_done,s_today,s_tmr,s_status,s_support) in enumerate(track_rows):
        tid = track_ids[tidx]
        track = get_track(tracks, tid)
        titems = [i for i in items if i.get('track') == tid]
        tdone  = [i for i in titems if i.get('status') in ['مكتملة','مكتمل','معتمدة','معتمد']]
        tactive= [i for i in titems if i.get('status') == 'قيد التنفيذ']
        tnext  = sorted([i for i in titems if i.get('due') and i.get('status') not in ['مكتملة','مكتمل','معتمدة','معتمد']], 
                        key=lambda x: x.get('due',''))
        trisks = [i for i in titems if i.get('type') in ['risks','مخاطرة'] and i.get('status') != 'مغلقة']
        
        set_slide_shape(s3, s_done,    tdone[0].get('title','—') if tdone else '—')
        set_slide_shape(s3, s_today,   tactive[0].get('title','—') if tactive else '—')
        set_slide_shape(s3, s_tmr,     tnext[0].get('title','—') if tnext else '—')
        set_slide_shape(s3, s_status,  status_label(track.get('status','')))
        set_slide_shape(s3, s_support, f"{len(trisks)} مخاطر مفتوحة" if trisks else 'لا يوجد')

    # Slide 5 — المخاطر والقرارات
    s5 = prs.slides[4]
    red_risks = [r for r in risks if r.get('status') in ['معرضة للخطر','معرض للخطر','At Risk']]
    yel_risks = [r for r in risks if r.get('status') in ['تحت المتابعة','قيد التنفيذ']]
    
    risk_rows_red = [
        ('Text 28','Text 22','Text 24','Text 26'),
    ]
    risk_rows_yel = [
        ('Text 38','Text 32','Text 34','Text 36'),
    ]
    risk_rows_grn = [
        ('Text 48','Text 42','Text 44','Text 46'),
    ]
    
    for (sq,mq,wq,tq), r in zip(risk_rows_red, red_risks):
        set_slide_shape(s5, sq, r.get('title','—'))
        set_slide_shape(s5, mq, r.get('owner','—'))
        set_slide_shape(s5, wq, format_date(r.get('due','')))
        set_slide_shape(s5, tq, 'تسريع الإجراءات')
    
    for (sq,mq,wq,tq), r in zip(risk_rows_yel, yel_risks):
        set_slide_shape(s5, sq, r.get('title','—'))
        set_slide_shape(s5, mq, r.get('owner','—'))
        set_slide_shape(s5, wq, format_date(r.get('due','')))
        set_slide_shape(s5, tq, 'متابعة')

    issue_detail = [
        ('Text 69','Text 67','Text 65','Text 63'),
        ('Text 77','Text 75','Text 73','Text 71'),
        ('Text 85','Text 83','Text 81','Text 79'),
    ]
    for (sq,dq,aq,iq), r in zip(issue_detail, risks[:3]):
        set_slide_shape(s5, sq, r.get('title','—'))
        set_slide_shape(s5, dq, r.get('title','—'))
        set_slide_shape(s5, aq, 'متابعة عاجلة')
        set_slide_shape(s5, iq, status_label(r.get('status','')))

    # Slide 6 — الجدول الزمني
    s6 = prs.slides[5]
    tasks = [i for i in items if i.get('type') not in ['risks','مخاطرة']]
    t_done   = [i for i in tasks if i.get('status') in ['مكتملة','مكتمل','معتمدة','معتمد']][:3]
    t_active = [i for i in tasks if i.get('status') == 'قيد التنفيذ'][:3]
    t_next   = sorted([i for i in tasks if i.get('due') and i.get('status') not in ['مكتملة','مكتمل','معتمدة','معتمد']], 
                      key=lambda x: x.get('due',''))[:3]
    
    fmt_list = lambda arr: '\n'.join(f"• {i.get('title','')}" for i in arr) or '• لا يوجد'
    set_slide_shape(s6, 'Text 42', fmt_list(t_done))
    set_slide_shape(s6, 'Text 31', fmt_list(t_active))
    set_slide_shape(s6, 'Text 20', fmt_list(t_next))

# ============================================================
# تعبئة تقرير مسار
# ============================================================
def fill_track(prs, state, track_id):
    tracks = state.get('tracks', [])
    items  = state.get('items', [])
    track  = get_track(tracks, track_id)
    titems = [i for i in items if i.get('track') == track_id]
    tdone  = [i for i in titems if i.get('status') in ['مكتملة','مكتمل','معتمدة','معتمد']]
    tactive= [i for i in titems if i.get('status') == 'قيد التنفيذ']
    tnext  = sorted([i for i in titems if i.get('due') and i.get('status') not in ['مكتملة','مكتمل','معتمدة','معتمد']],
                    key=lambda x: x.get('due',''))
    risks  = [i for i in titems if i.get('type') in ['risks','مخاطرة'] and i.get('status') != 'مغلقة']
    red_r  = [r for r in risks if r.get('status') in ['معرضة للخطر','معرض للخطر','At Risk']]
    yel_r  = [r for r in risks if r.get('status') in ['تحت المتابعة','قيد التنفيذ']]

    # Slide 1 — غلاف (ثابت)
    
    # Slide 2 — الملخص
    s2 = prs.slides[1]
    set_slide_shape(s2, 'Text 12', f"الحالة المختارة: {status_label(track.get('status',''))}\nنسبة إنجاز اليوم: {track.get('progress',0)}٪")
    
    done_t = [i.get('title','') for i in tdone[:3]]
    while len(done_t) < 3: done_t.append('—')
    set_slide_shape(s2, 'Text 15', '\n'.join(done_t))
    set_slide_shape(s2, 'Text 18', risks[0].get('title','لا يوجد') + ('\nآخر موعد: ' + format_date(risks[0].get('due',''))) if risks else 'لا يوجد')
    
    # جدول التحديث التفصيلي — 4 صفوف
    detail_rows = [
        ('Text 44','Text 40','Text 38','Text 36','Text 34'),  # أنشطة رئيسية
        ('Text 56','Text 52','Text 50','Text 48','Text 46'),  # مخرجات/اعتمادات
        ('Text 68','Text 64','Text 62','Text 60','Text 58'),  # تنسيق مع جهات
        ('Text 80','Text 76','Text 74','Text 72','Text 70'),  # موردين/أطراف
    ]
    all_items = tdone + tactive + tnext
    for row_idx, (_, s_done, s_active, s_planned, _) in enumerate(detail_rows):
        base = row_idx * 2
        i_done    = all_items[base] if base < len(all_items) else None
        i_active  = all_items[base+1] if base+1 < len(all_items) else None
        set_slide_shape(s2, s_done,    i_done.get('title','—') if i_done else '—')
        set_slide_shape(s2, s_active,  i_active.get('title','—') if i_active else '—')
        set_slide_shape(s2, s_planned, '—')

    # Slide 3 — الأنشطة التفصيلية (8 صفوف)
    s3 = prs.slides[2]
    # كل صف: النشاط(35+16n), المالك(33+16n), وقت(31+16n), حالة(27+16n), نسبة(25+16n)
    activity_shapes = [
        ('Text 35','Text 33','Text 31','Text 27','Text 25'),
        ('Text 51','Text 49','Text 47','Text 43','Text 41'),
        ('Text 67','Text 65','Text 63','Text 59','Text 57'),
        ('Text 83','Text 81','Text 79','Text 75','Text 73'),
        ('Text 99','Text 97','Text 95','Text 91','Text 89'),
        ('Text 115','Text 113','Text 111','Text 107','Text 105'),
        ('Text 131','Text 129','Text 127','Text 123','Text 121'),
        ('Text 147','Text 145','Text 143','Text 139','Text 137'),
    ]
    display_items = (tdone + tactive + tnext)[:8]
    for idx, (t_title, t_owner, t_due, t_status, t_pct) in enumerate(activity_shapes):
        if idx < len(display_items):
            item = display_items[idx]
            set_slide_shape(s3, t_title,  item.get('title','—'))
            set_slide_shape(s3, t_owner,  item.get('owner','—'))
            set_slide_shape(s3, t_due,    format_date(item.get('due','')))
            set_slide_shape(s3, t_status, status_label(item.get('status','')))
            set_slide_shape(s3, t_pct,    f"{item.get('progress',0)}٪")
        else:
            set_slide_shape(s3, t_title,  '—')
            set_slide_shape(s3, t_owner,  '')
            set_slide_shape(s3, t_due,    '')
            set_slide_shape(s3, t_status, '')
            set_slide_shape(s3, t_pct,    '')

    # Slide 4 — المخاطر
    s4 = prs.slides[3]
    risk_map = [
        (red_r, [('Text 28','Text 17','Text 19','Text 21','Text 23')]),
        (yel_r, [('Text 38','Text 29','Text 31','Text 33','Text 35')]),
        (risks, [('Text 48','Text 41','Text 43','Text 45','Text 47')]),
    ]
    for risk_list, shape_rows in risk_map:
        for (sq, mq, wq, tq, dq), r in zip(shape_rows, risk_list):
            set_slide_shape(s4, sq, r.get('title','—'))
            set_slide_shape(s4, mq, r.get('owner','—') if mq != 'Text 17' else r.get('owner','—'))
            set_slide_shape(s4, wq, format_date(r.get('due','')))
            set_slide_shape(s4, tq, 'تسريع الإجراءات')

    issue_rows = [
        ('Text 69','Text 67','Text 65','Text 63'),
        ('Text 77','Text 75','Text 73','Text 71'),
        ('Text 85','Text 83','Text 81','Text 79'),
    ]
    for (sq,dq,aq,iq), r in zip(issue_rows, risks[:3]):
        set_slide_shape(s4, sq, r.get('title','—'))
        set_slide_shape(s4, dq, r.get('title','—'))
        set_slide_shape(s4, aq, 'متابعة عاجلة')
        set_slide_shape(s4, iq, status_label(r.get('status','')))

    # Slide 6 — الجدول الزمني
    s6 = prs.slides[5]
    fmt_list = lambda arr: '\n'.join(f"• {i.get('title','')}" for i in arr) or '• لا يوجد'
    set_slide_shape(s6, 'Text 33', fmt_list(tdone[:3]))
    set_slide_shape(s6, 'Text 24', fmt_list(tactive[:3]))
    set_slide_shape(s6, 'Text 15', fmt_list(tnext[:3]))

# ============================================================
# Main
# ============================================================
def main():
    data = json.loads(sys.stdin.read())
    report_type = data['type']
    state = data['state']
    tpl_dir = data.get('tpl_dir', 'templates')
    
    tpl_map = {
        'comprehensive': 'comprehensive.pptx',
        'أ': 'track-a.pptx',
        'ب': 'track-b.pptx',
        'ج': 'track-c.pptx',
        'د': 'track-d.pptx',
    }
    
    tpl_file = tpl_map.get(report_type)
    if not tpl_file:
        sys.stderr.write(f"نوع تقرير غير معروف: {report_type}\n")
        sys.exit(1)
    
    tpl_path = os.path.join(tpl_dir, tpl_file)
    prs = Presentation(tpl_path)
    
    if report_type == 'comprehensive':
        fill_comprehensive(prs, state)
    else:
        fill_track(prs, state, report_type)
    
    buf = io.BytesIO()
    prs.save(buf)
    sys.stdout.buffer.write(buf.getvalue())

if __name__ == '__main__':
    main()
