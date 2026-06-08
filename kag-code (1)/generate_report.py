#!/usr/bin/env python3
"""
generate_report.py v3 — مولّد تقارير PPTX
يستخدم القوالب الأصلية ويعبّئ البيانات في أماكنها الصحيحة بالضبط
"""
import sys, json, io, os
from datetime import datetime
from pptx import Presentation

TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")

ORIGINALS = {
    "comprehensive": "comprehensive.pptx",
    "أ": "track-a.pptx",
    "ب": "track-b.pptx",
    "ج": "track-c.pptx",
    "د": "track-d.pptx",
}

# ============================================================
# مساعدات
# ============================================================
def today_str():
    return datetime.now().strftime("%Y-%m-%d")

def week_num():
    return datetime.now().isocalendar()[1]

def fmt_date(d):
    if not d: return "—"
    try: return datetime.strptime(str(d)[:10], "%Y-%m-%d").strftime("%Y/%m/%d")
    except: return str(d)[:10]

def status_ar(s):
    if not s: return "—"
    if s in ["مكتملة","مكتمل","معتمدة","معتمد","ضمن المسار"]: return "أخضر ✓"
    if s in ["قيد التنفيذ","تحت المتابعة"]: return "أصفر"
    if s in ["معرضة للخطر","معرض للخطر","متأخرة"]: return "أحمر ✗"
    return s

def is_done(i): return i.get("status","") in ["مكتملة","مكتمل","معتمدة","معتمد"]
def is_active(i): return i.get("status","") == "قيد التنفيذ"
def is_risk(i): return i.get("type","") in ["مخاطرة","مخاطر","risks"]

def get(lst, idx, key="title", default="—"):
    return lst[idx].get(key, default) if idx < len(lst) else default

def bullet_lines(lst, key="title"):
    if not lst: return "لا يوجد"
    return "\n".join(f"• {i.get(key,'')}" for i in lst)

def F(slide, name, text):
    """ابحث عن shape باسمه وعيّن نصه مع الحفاظ على التنسيق"""
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            tf = shape.text_frame
            if not tf.paragraphs: return
            # احذف paragraphs الزائدة
            while len(tf.paragraphs) > 1:
                tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
            para = tf.paragraphs[0]
            # احذف runs الزائدة
            while len(para.runs) > 1:
                para.runs[-1]._r.getparent().remove(para.runs[-1]._r)
            # عيّن النص
            t = str(text) if text else "—"
            if para.runs:
                para.runs[0].text = t
            else:
                para.text = t
            return

# ============================================================
# التقرير الشامل
# ============================================================
def fill_comprehensive(prs, state):
    tracks = state.get("tracks", [])
    items  = state.get("items", [])
    risks  = [i for i in items if is_risk(i) and i.get("status") != "مغلقة"]
    done   = [i for i in items if is_done(i)]
    active = [i for i in items if is_active(i)]

    overall = round(sum(t.get("progress",0) for t in tracks)/len(tracks)) if tracks else 0
    worst = "أحمر" if any(t.get("status","") in ["معرض للخطر","معرضة للخطر"] for t in tracks) \
        else "أصفر" if any(t.get("status","") == "تحت المتابعة" for t in tracks) else "أخضر"

    def gt(tid): return next((t for t in tracks if t.get("track")==tid), {})

    # شريحة 1 — الغلاف
    s1 = prs.slides[0]
    F(s1, "Text 4", f"الأسبوع {week_num()} ٢٠٢٦  |  {today_str()}")

    # شريحة 2 — الملخص التنفيذي
    s2 = prs.slides[1]
    F(s2, "Text 17", f"الحالة: {worst}  |  نسبة إنجاز اليوم: {overall}٪")
    F(s2, "Text 21", bullet_lines(done[:3]))
    F(s2, "Text 25", bullet_lines(active[:3]))
    # القضايا الحرجة — 3 صفوف
    risk_rows = [
        ("Text 44","Text 46","Text 48","Text 50","Shape 51"),
        ("Text 54","Text 56","Text 58","Text 60","Text 62"),
        ("Text 64","Text 66","Text 68","Text 70","Text 72"),
    ]
    for idx, (st, sm, sa, sar, sq) in enumerate(risk_rows):
        r = risks[idx] if idx < len(risks) else None
        F(s2, st,  fmt_date(r.get("due",""))  if r else "—")
        F(s2, sm,  r.get("owner","—")         if r else "—")
        F(s2, sa,  "متابعة عاجلة"              if r else "—")
        F(s2, sar, "تأثير على الجدول"          if r else "—")
        F(s2, sq,  r.get("title","—")          if r else "—")

    # شريحة 3 — المسارات الأربعة
    s3 = prs.slides[2]
    track_map = [
        ("أ","Text 32","Text 30","Text 28","Text 26","Text 24"),
        ("ب","Text 44","Text 42","Text 40","Text 38","Text 36"),
        ("ج","Text 56","Text 54","Text 52","Text 50","Text 48"),
        ("د","Text 68","Text 66","Text 64","Text 62","Text 60"),
    ]
    for tid, amsl, alyom, ghad, hal, da3m in track_map:
        t = gt(tid)
        ti = [i for i in items if i.get("track")==tid]
        td = [i for i in ti if is_done(i)]
        ta = [i for i in ti if is_active(i)]
        tn = sorted([i for i in ti if not is_risk(i) and i.get("due","")>today_str()], key=lambda x:x.get("due",""))
        tr = [i for i in ti if is_risk(i) and i.get("status")!="مغلقة"]
        F(s3, amsl,  get(td,0) if td else "—")
        F(s3, alyom, get(ta,0) if ta else "—")
        F(s3, ghad,  get(tn,0) if tn else "—")
        F(s3, hal,   status_ar(t.get("status","")))
        F(s3, da3m,  f"{len(tr)} مخاطر" if tr else "لا يوجد")

    # شريحة 4 — السلامة (نتركها كما هي)

    # شريحة 5 — المخاطر والقرارات
    s5 = prs.slides[4]
    red = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
    yel = [r for r in risks if r.get("status","") in ["تحت المتابعة","قيد التنفيذ"]]
    grn = [r for r in risks if r not in red and r not in yel]
    buckets = [
        (red, "Text 28","Text 22","Text 24","Text 26","Text 30"),
        (yel, "Text 38","Text 32","Text 34","Text 36","Text 40"),
        (grn, "Text 48","Text 42","Text 44","Text 46","Text 50"),
    ]
    for lst, sq, sm, st, stow, _ in buckets:
        r = lst[0] if lst else None
        F(s5, sq,   r.get("title","—")       if r else "—")
        F(s5, sm,   r.get("owner","—")        if r else "—")
        F(s5, st,   fmt_date(r.get("due","")) if r else "—")
        F(s5, stow, "متابعة عاجلة"            if r else "—")
    dec_rows = [
        ("Text 69","Text 67","Text 65","Text 63"),
        ("Text 77","Text 75","Text 73","Text 71"),
        ("Text 85","Text 83","Text 81","Text 79"),
    ]
    for idx, (sw, sa, sar, sr) in enumerate(dec_rows):
        r = risks[idx] if idx < len(risks) else None
        F(s5, sw,  r.get("title","—") if r else "—")
        F(s5, sa,  "تأثير على الجدول" if r else "—")
        F(s5, sar, "متابعة عاجلة"     if r else "—")
        F(s5, sr,  f"خطر-{idx+1}"     if r else "—")

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    all_tasks = [i for i in items if not is_risk(i)]
    td = [i for i in all_tasks if is_done(i)][:4]
    ta = [i for i in all_tasks if is_active(i)][:4]
    tn = sorted([i for i in all_tasks if i.get("due","")>today_str()], key=lambda x:x.get("due",""))[:4]
    F(s6, "Text 42", bullet_lines(td))
    F(s6, "Text 31", bullet_lines(ta))
    F(s6, "Text 20", bullet_lines(tn))

# ============================================================
# تقارير المسارات (أ، ب، ج، د — نفس البنية)
# ============================================================
def fill_track(prs, track_key, state):
    tracks = state.get("tracks", [])
    items  = state.get("items", [])
    track  = next((t for t in tracks if t.get("track")==track_key), {})
    ti     = [i for i in items if i.get("track")==track_key]
    risks  = [i for i in ti if is_risk(i) and i.get("status")!="مغلقة"]
    done   = [i for i in ti if is_done(i)]
    active = [i for i in ti if is_active(i)]
    tasks  = [i for i in ti if not is_risk(i)]
    upcoming = sorted([i for i in ti if not is_risk(i) and i.get("due","")>today_str()], key=lambda x:x.get("due",""))
    progress = track.get("progress", 0)

    # شريحة 1 — الغلاف: لا تعديل (اسم المسار موجود في القالب)

    # شريحة 2 — الملخص اليومي
    s2 = prs.slides[1]
    F(s2, "Text 12", f"الحالة المختارة: {status_ar(track.get('status',''))}  |  نسبة إنجاز اليوم: {progress}٪")
    # ملخص الإنجاز — 3 عناصر
    F(s2, "Text 15",
      f"{get(done,0)}\n{get(done,1)}\n{get(done,2)}")
    # الدعم المطلوب
    F(s2, "Text 18",
      f"{get(risks,0)}  |  آخر موعد: {fmt_date(risks[0].get('due','')) if risks else '—'}"
      if risks else "لا يوجد دعم عاجل مطلوب")
    # التحديث التفصيلي — 4 صفوف × 3 أعمدة
    detail_rows = [
        ("Text 42","Text 40","Text 38", 0),  # أنشطة رئيسية
        ("Text 54","Text 52","Text 50", 1),  # مخرجات/اعتمادات
        ("Text 66","Text 64","Text 62", 2),  # تنسيق مع جهات
        ("Text 78","Text 76","Text 74", 3),  # موردين/أطراف
    ]
    for amsl, alyom, ghad, idx in detail_rows:
        F(s2, amsl,  get(done,    idx) if idx < len(done)     else "—")
        F(s2, alyom, get(active,  idx) if idx < len(active)   else "—")
        F(s2, ghad,  get(upcoming,idx) if idx < len(upcoming) else "—")

    # شريحة 3 — تفاصيل الأنشطة (8 صفوف)
    s3 = prs.slides[2]
    task_rows = [
        ("Text 35","Text 33","Text 31","Text 29","Text 27","Text 25"),
        ("Text 51","Text 49","Text 47","Text 45","Text 43","Text 41"),
        ("Text 67","Text 65","Text 63","Text 61","Text 59","Text 57"),
        ("Text 83","Text 81","Text 79","Text 77","Text 75","Text 73"),
        ("Text 99","Text 97","Text 95","Text 93","Text 91","Text 89"),
        ("Text 115","Text 113","Text 111","Text 109","Text 107","Text 105"),
        ("Text 131","Text 129","Text 127","Text 125","Text 123","Text 121"),
        ("Text 147","Text 145","Text 143","Text 141","Text 139","Text 137"),
    ]
    for idx, (stitle, sowner, sstart, send, sstatus, spct) in enumerate(task_rows):
        t = tasks[idx] if idx < len(tasks) else None
        F(s3, stitle,  t.get("title","—")             if t else "—")
        F(s3, sowner,  t.get("owner","—")              if t else "—")
        F(s3, sstart,  fmt_date(t.get("due",""))       if t else "—")
        F(s3, send,    fmt_date(t.get("due",""))        if t else "—")
        F(s3, sstatus, status_ar(t.get("status",""))   if t else "—")
        F(s3, spct,    f"{t.get('progress',0)}٪"        if t else "—")

    # شريحة 4 — المخاطر والقرارات
    s4 = prs.slides[3]
    red = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
    yel = [r for r in risks if r.get("status","") in ["تحت المتابعة","قيد التنفيذ"]]
    grn = [r for r in risks if r not in red and r not in yel]
    risk_buckets = [
        (red, "Text 25","Text 19","Text 21","Text 23","Text 17"),
        (yel, "Text 37","Text 29","Text 31","Text 33","Text 35"),  # fixed: أص row
        (grn, "Text 49","Text 41","Text 43","Text 45","Text 47"),  # fixed: أ row
    ]
    for lst, sq, stime1, sm, stow, stime2 in risk_buckets:
        r = lst[0] if lst else None
        F(s4, sq,     r.get("title","—")       if r else "—")
        F(s4, sm,     r.get("owner","—")        if r else "—")
        F(s4, stow,   "متابعة عاجلة"            if r else "—")
        F(s4, stime1, fmt_date(r.get("due","")) if r else "—")
        F(s4, stime2, fmt_date(r.get("due","")) if r else "—")
    dec_rows4 = [
        ("Text 71","Text 69","Text 67","Text 65"),
        ("Text 79","Text 77","Text 75","Text 73"),
        ("Text 87","Text 85","Text 83","Text 81"),
    ]
    for idx, (sw, sa, sar, sr) in enumerate(dec_rows4):
        r = risks[idx] if idx < len(risks) else None
        F(s4, sw,  r.get("title","—") if r else "—")
        F(s4, sa,  "تأثير على الجدول" if r else "—")
        F(s4, sar, "متابعة"            if r else "—")
        F(s4, sr,  f"خطر-{idx+1}"     if r else "—")

    # شريحة 5 — السلامة: نتركها كما هي

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    td6 = [i for i in tasks if is_done(i)][:4]
    ta6 = [i for i in tasks if is_active(i)][:4]
    tn6 = upcoming[:4]
    F(s6, "Text 33", bullet_lines(td6))
    F(s6, "Text 24", bullet_lines(ta6))
    F(s6, "Text 15", bullet_lines(tn6))

    # شريحة 7 — التوصيات: نتركها فارغة للتعبئة اليدوية

# ============================================================
# الدالة الرئيسية
# ============================================================
def generate_report(report_type, state):
    tpl_file = ORIGINALS.get(report_type)
    if not tpl_file:
        raise ValueError(f"نوع تقرير غير معروف: {report_type}")
    tpl_path = os.path.join(TEMPLATES_DIR, tpl_file)
    if not os.path.exists(tpl_path):
        raise FileNotFoundError(f"القالب غير موجود: {tpl_path}")
    prs = Presentation(tpl_path)
    if report_type == "comprehensive":
        fill_comprehensive(prs, state)
    else:
        fill_track(prs, report_type, state)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

if __name__ == "__main__":
    data  = json.loads(sys.stdin.read())
    result = generate_report(data.get("type","comprehensive"), data.get("state",{}))
    sys.stdout.buffer.write(result)
