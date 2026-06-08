# -*- coding: utf-8 -*-
"""محرّك التقرير اليومي الشامل — يقرأ التحديث اليومي من CSV ويعبّي قالب العميل."""
import csv
from pptx import Presentation
from pptx.util import Emu

TEMPLATE="daily_comprehensive.pptx"; OUT="out_comprehensive.pptx"
DAILY_CSV="daily_update_template.csv"; ISSUES_CSV="critical_issues_template.csv"

def read_daily(p):
    d={}
    with open(p,encoding="utf-8-sig") as f:
        for r in csv.DictReader(f): d[r["القسم"].strip()]=r
    return d
def read_issues(p):
    o=[]
    with open(p,encoding="utf-8-sig") as f:
        for r in csv.DictReader(f): o.append([r["الأثر"].strip(),r["الإجراء"].strip(),r["المالك"].strip(),r["الموعد"].strip()])
    return o
def split_multi(v,n=3):
    parts=[x.strip() for x in (v or "").split("|") if x.strip()]
    return (parts+[""]*n)[:n]

daily=read_daily(DAILY_CSV); issues=read_issues(ISSUES_CSV); g=daily.get("عام",{})
DATA={"status":g.get("الحالة","-"),"daily_pct":g.get("نسبة_الإنجاز","-"),
 "yesterday":split_multi(g.get("إنجازات_أمس","")),"today":split_multi(g.get("أنشطة_اليوم","")),
 "notes":g.get("ملاحظات",""),"issues":issues}

def set_lines(sh,lines):
    for i,p in enumerate(sh.text_frame.paragraphs):
        val=lines[i] if i<len(lines) else ""
        if p.runs:
            p.runs[0].text=val
            for r in p.runs[1:]: r.text=""
        else: p.text=val
def set_text(sh,v): set_lines(sh,[v])
def text_of(sh): return sh.text_frame.text.strip() if sh.has_text_frame else ""

prs=Presentation(TEMPLATE); S=prs.slides
for sh in S[0].shapes:
    if sh.has_text_frame and "رقم الأسبوع" in text_of(sh):
        set_text(sh,"التقرير اليومي الشامل — مشروع افتتاح حدائق الملك عبدالله")
s2=S[1]
for sh in s2.shapes:
    if not sh.has_text_frame: continue
    t=text_of(sh)
    if t.startswith("الحالة:"): set_lines(sh,["الحالة: "+DATA["status"],"نسبة إنجاز اليوم: "+DATA["daily_pct"]+"٪"])
    elif t.replace("\n"," ").startswith("١ إنجاز"): set_lines(sh,DATA["yesterday"])
    elif t.replace("\n"," ").startswith("١ نشاط"): set_lines(sh,DATA["today"])
rows_y=[3.95,4.37,4.79]; issue_boxes=[]
for sh in s2.shapes:
    if sh.has_text_frame and text_of(sh) in ("قضية","أثر","إجراء","مالك","وقت"):
        issue_boxes.append((round(Emu(sh.top).inches,2),text_of(sh),sh))
for ri,y in enumerate(rows_y):
    if ri>=len(DATA["issues"]): break
    impact,action,owner,due=DATA["issues"][ri]
    valby={"أثر":impact,"إجراء":action,"مالك":owner,"وقت":due,"قضية":"قضية "+str(ri+1)}
    for by,label,sh in issue_boxes:
        if abs(by-y)<=0.3 and label in valby: set_text(sh,valby[label])
cand=[]
for sh in s2.shapes:
    if sh.has_text_frame and not text_of(sh):
        try: x,y,h=Emu(sh.left).inches,Emu(sh.top).inches,Emu(sh.height).inches
        except: continue
        if x<6 and y>3.0 and h>1.0: cand.append((h,sh))
# (ملاحظات اليوم: يُضبط استهدافها بدقة في التكرار القادم)
pass
prs.save(OUT)
print("saved",OUT,"| الحالة:",DATA["status"],"| النسبة:",DATA["daily_pct"],"| قضايا:",len(DATA["issues"]))
