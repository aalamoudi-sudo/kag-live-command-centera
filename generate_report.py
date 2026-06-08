#!/usr/bin/env python3
"""
generate_report.py — مولّد تقارير PPTX
يفتح القوالب الأصلية ويعدّل النصوص فقط مع الحفاظ الكامل على التصميم
"""
import sys, json, io, os
from datetime import datetime
from pptx import Presentation

TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")

TEMPLATE_FILES = {
    "comprehensive": "comprehensive.pptx",
    "أ": "track-a.pptx",
    "ب": "track-b.pptx",
    "ج": "track-c.pptx",
    "د": "track-d.pptx",
}

# ============================================================
# مساعدات
# ============================================================
def today_str():
    return datetime.now().strftime("%Y-%m-%d")

def week_num():
    return datetime.now().isocalendar()[1]

def fmt_date(d):
    if not d: return "—"
    try:
        return datetime.strptime(str(d)[:10], "%Y-%m-%d").strftime("%Y/%m/%d")
    except:
        return str(d)

def status_ar(s):
    if not s: return "—"
    if s in ["مكتملة","مكتمل","معتمدة","معتمد","ضمن المسار"]: return "أخضر ✓"
    if s in ["قيد التنفيذ","تحت المتابعة"]: return "أصفر"
    if s in ["معرضة للخطر","معرض للخطر","متأخرة"]: return "أحمر ✗"
    return s

def is_done(i):
    return i.get("status","") in ["مكتملة","مكتمل","معتمدة","معتمد"]

def is_active(i):
    return i.get("status","") == "قيد التنفيذ"

def is_risk(i):
    return i.get("type","") in ["مخاطرة","مخاطر","risks"]

def safe_set(shape, text):
    """تعديل نص shape مع الحفاظ الكامل على التنسيق"""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    # احذف كل paragraphs ما عدا الأول
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)
    para = tf.paragraphs[0]
    # احذف كل runs ما عدا الأول
    while len(para.runs) > 1:
        para.runs[-1]._r.getparent().remove(para.runs[-1]._r)
    # عيّن النص
    if para.runs:
        para.runs[0].text = str(text) if text else "—"
    else:
        para.text = str(text) if text else "—"

def find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def F(slide, name, text):
    """اختصار: ابحث وعيّن"""
    safe_set(find(slide, name), text)

def title_items(items, track=None):
    """استخرج عناوين المهام مرتبة"""
    filtered = [i for i in items if not is_risk(i)]
    if track:
        filtered = [i for i in filtered if i.get("track") == track]
    done_list   = [i.get("title","") for i in filtered if is_done(i)][:4]
    active_list = [i.get("title","") for i in filtered if is_active(i)][:4]
    upcoming    = sorted(
        [i for i in filtered if not is_done(i) and i.get("due","") > today_str()],
        key=lambda x: x.get("due","")
    )[:4]
    return done_list, active_list, [i.get("title","") for i in upcoming]

def bullet(lst):
    return "  •  ".join(lst) if lst else "لا يوجد"

# ============================================================
# التقرير الشامل
# ============================================================
def fill_comprehensive(prs, state):
    tracks = state.get("tracks", [])
    items  = state.get("items", [])
    risks  = [i for i in items if is_risk(i) and i.get("status") != "مغلقة"]
    
    overall = round(sum(t.get("progress",0) for t in tracks)/len(tracks)) if tracks else 0
    worst = "أحمر" if any(t.get("status","") in ["معرض للخطر","معرضة للخطر"] for t in tracks) \
        else "أصفر" if any(t.get("status","") == "تحت المتابعة" for t in tracks) else "أخضر"

    def gt(tid): return next((t for t in tracks if t.get("track")==tid), {})

    # شريحة 1 — الغلاف
    s1 = prs.slides[0]
    F(s1, "Text 4", f"الأسبوع {week_num()} ٢٠٢٦  |  {today_str()}")

    # شريحة 2 — الملخص التنفيذي
    s2 = prs.slides[1]
    done_l, active_l, _ = title_items(items)
    F(s2, "Text 17", f"الحالة: {worst}  |  نسبة إنجاز اليوم: {overall}٪")
    F(s2, "Text 21", bullet(done_l[:3]))
    F(s2, "Text 25", bullet(active_l[:3]))
    # القضايا الحرجة
    risk_shapes = [
        ("Text 44","Text 46","Text 48","Text 50","Shape 51"),
        ("Text 54","Text 56","Text 58","Text 60","Text 62"),
        ("Text 64","Text 66","Text 68","Text 70","Text 72"),
    ]
    for idx, (st,sm,sa,sar,sq) in enumerate(risk_shapes):
        r = risks[idx] if idx < len(risks) else None
        F(s2, st,  fmt_date(r.get("due","")) if r else "—")
        F(s2, sm,  r.get("owner","—")        if r else "—")
        F(s2, sa,  "متابعة عاجلة"             if r else "—")
        F(s2, sar, "تأثير على الجدول"         if r else "—")
        F(s2, sq,  r.get("title","—")         if r else "—")

    # شريحة 3 — المسارات
    s3 = prs.slides[2]
    # كل مسار = 5 أعمدة: إنجاز أمس، خطة اليوم، خطة الغد، الحالة، الدعم
    track_rows = [
        ("أ", "Text 32","Text 30","Text 28","Text 26","Text 24"),
        ("ب", "Text 44","Text 42","Text 40","Text 38","Text 36"),
        ("ج", "Text 56","Text 54","Text 52","Text 50","Text 48"),
        ("د", "Text 68","Text 66","Text 64","Text 62","Text 60"),
    ]
    for tid, amsl, alyom, ghad, hal, da3m in track_rows:
        t = gt(tid)
        dl, al, ul = title_items(items, tid)
        tr = [i for i in items if i.get("track")==tid and is_risk(i) and i.get("status")!="مغلقة"]
        F(s3, amsl,  dl[0] if dl else "—")
        F(s3, alyom, al[0] if al else "—")
        F(s3, ghad,  ul[0] if ul else "—")
        F(s3, hal,   status_ar(t.get("status","")))
        F(s3, da3m,  f"{len(tr)} مخاطر مفتوحة" if tr else "لا يوجد")

    # شريحة 4 — السلامة
    s4 = prs.slides[3]
    crit = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
    F(s4, "Text 67", crit[0].get("title","—") if crit else "—")
    F(s4, "Text 69", crit[0].get("owner","—") if crit else "—")
    F(s4, "Text 79", crit[1].get("title","—") if len(crit)>1 else "—")
    F(s4, "Text 81", crit[1].get("owner","—") if len(crit)>1 else "—")

    # شريحة 5 — المخاطر والقرارات
    s5 = prs.slides[4]
    red  = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
    yel  = [r for r in risks if r.get("status","") in ["تحت المتابعة","قيد التنفيذ"]]
    grn  = [r for r in risks if r not in red and r not in yel]
    bucket_shapes = [
        ("ح",  red,  "Text 28","Text 22","Text 24","Text 26"),
        ("أص", yel,  "Text 38","Text 32","Text 34","Text 36"),
        ("أ",  grn,  "Text 48","Text 42","Text 44","Text 46"),
    ]
    for _,lst,sq,sm,st,stow in bucket_shapes:
        r = lst[0] if lst else None
        F(s5, sq,   r.get("title","—")       if r else "—")
        F(s5, sm,   r.get("owner","—")        if r else "—")
        F(s5, st,   fmt_date(r.get("due","")) if r else "—")
        F(s5, stow, "متابعة عاجلة"            if r else "—")
    dec_shapes = [
        ("Text 69","Text 67","Text 65","Text 63"),
        ("Text 77","Text 75","Text 73","Text 71"),
        ("Text 85","Text 83","Text 81","Text 79"),
    ]
    for idx,(sw,sa,sar,sr) in enumerate(dec_shapes):
        r = risks[idx] if idx < len(risks) else None
        F(s5, sw, r.get("title","—")       if r else "—")
        F(s5, sa, "تأثير على الجدول"        if r else "—")
        F(s5, sar,"متابعة عاجلة"            if r else "—")
        F(s5, sr, f"خطر-{idx+1}"            if r else "—")

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    dl, al, ul = title_items(items)
    F(s6, "Text 42", bullet(dl[:4]))
    F(s6, "Text 31", bullet(al[:4]))
    F(s6, "Text 20", bullet(ul[:4]))

# ============================================================
# تقرير المسار
# ============================================================
def fill_track(prs, track_key, state):
    tracks = state.get("tracks", [])
    items  = state.get("items", [])
    track  = next((t for t in tracks if t.get("track")==track_key), {})
    ti     = [i for i in items if i.get("track")==track_key]
    risks  = [i for i in ti if is_risk(i) and i.get("status")!="مغلقة"]
    done   = [i for i in ti if is_done(i)]
    active = [i for i in ti if is_active(i)]
    tasks  = [i for i in ti if not is_risk(i)]
    upcoming = sorted([i for i in ti if not is_risk(i) and i.get("due","")>today_str()], key=lambda x:x.get("due",""))
    progress = track.get("progress", 0)

    # شريحة 1 — الغلاف: لا تعديل

    # شريحة 2 — الملخص
    s2 = prs.slides[1]
    F(s2, "Text 12", f"الحالة: {status_ar(track.get('status',''))}  |  نسبة إنجاز اليوم: {progress}٪")
    F(s2, "Text 15",
      (done[0].get("title","—") if done else "—") + "  |  " +
      (done[1].get("title","—") if len(done)>1 else "—") + "  |  " +
      (done[2].get("title","—") if len(done)>2 else "—"))
    F(s2, "Text 18",
      (risks[0].get("title","") + "  |  " + fmt_date(risks[0].get("due",""))) if risks else "لا يوجد دعم عاجل")
    # صفوف التحديث التفصيلي
    row_data = [
        ("Text 42","Text 40","Text 38", done, active, upcoming, 0),
        ("Text 54","Text 52","Text 50", done, active, upcoming, 1),
        ("Text 66","Text 64","Text 62", done, active, upcoming, 2),
        ("Text 78","Text 76","Text 74", done, active, upcoming, 3),
    ]
    for amsl, alyom, ghad, dl, al, ul, idx in row_data:
        F(s2, amsl,  dl[idx].get("title","—") if idx<len(dl) else "—")
        F(s2, alyom, al[idx].get("title","—") if idx<len(al) else "—")
        F(s2, ghad,  ul[idx].get("title","—") if idx<len(ul) else "—")

    # شريحة 3 — تفاصيل الأنشطة
    s3 = prs.slides[2]
    task_rows = [
        ("Text 35","Text 33","Text 31","Text 29","Text 27","Text 25"),
        ("Text 51","Text 49","Text 47","Text 45","Text 43","Text 41"),
        ("Text 67","Text 65","Text 63","Text 61","Text 59","Text 57"),
        ("Text 83","Text 81","Text 79","Text 77","Text 75","Text 73"),
        ("Text 99","Text 97","Text 95","Text 93","Text 91","Text 89"),
        ("Text 115","Text 113","Text 111","Text 109","Text 107","Text 105"),
        ("Text 131","Text 129","Text 127","Text 125","Text 123","Text 121"),
        ("Text 147","Text 145","Text 143","Text 141","Text 139","Text 137"),
    ]
    for idx, (st,sown,ss,se,shal,spct) in enumerate(task_rows):
        t = tasks[idx] if idx < len(tasks) else None
        F(s3, st,   t.get("title","—")                if t else "—")
        F(s3, sown, t.get("owner","—")                if t else "—")
        F(s3, ss,   fmt_date(t.get("due",""))          if t else "—")
        F(s3, se,   fmt_date(t.get("due",""))          if t else "—")
        F(s3, shal, status_ar(t.get("status",""))      if t else "—")
        F(s3, spct, f"{t.get('progress',0)}٪"          if t else "—")

    # شريحة 4 — المخاطر
    s4 = prs.slides[3]
    red = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
    yel = [r for r in risks if r.get("status","") in ["تحت المتابعة","قيد التنفيذ"]]
    grn = [r for r in risks if r not in red and r not in yel]
    b4 = [
        (red, "Text 25","Text 21","Text 23","Text 17","Text 19"),
        (yel, "Text 37","Text 33","Text 35","Text 29","Text 31"),
        (grn, "Text 49","Text 45","Text 47","Text 41","Text 43"),
    ]
    for lst,sq,sm,st,stow,stime in b4:
        r = lst[0] if lst else None
        F(s4, sq,    r.get("title","—")       if r else "—")
        F(s4, sm,    r.get("owner","—")        if r else "—")
        F(s4, st,    fmt_date(r.get("due","")) if r else "—")
        F(s4, stow,  "متابعة عاجلة"            if r else "—")
        F(s4, stime, fmt_date(r.get("due","")) if r else "—")
    dec4 = [
        ("Text 71","Text 69","Text 67","Text 65"),
        ("Text 79","Text 77","Text 75","Text 73"),
        ("Text 87","Text 85","Text 83","Text 81"),
    ]
    for idx,(sw,sa,sar,sr) in enumerate(dec4):
        r = risks[idx] if idx < len(risks) else None
        F(s4, sw,  r.get("title","—")       if r else "—")
        F(s4, sa,  "تأثير على الجدول"        if r else "—")
        F(s4, sar, "متابعة"                  if r else "—")
        F(s4, sr,  f"خطر-{idx+1}"            if r else "—")

    # شريحة 5 — السلامة: نتركها كما هي

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    dl = [i.get("title","") for i in tasks if is_done(i)][:4]
    al = [i.get("title","") for i in tasks if is_active(i)][:4]
    ul = [i.get("title","") for i in upcoming][:4]
    F(s6, "Text 33", bullet(dl))
    F(s6, "Text 24", bullet(al))
    F(s6, "Text 15", bullet(ul))

# ============================================================
# الدالة الرئيسية
# ============================================================
def generate_report(report_type, state):
    tpl_file = TEMPLATE_FILES.get(report_type)
    if not tpl_file:
        raise ValueError(f"نوع تقرير غير معروف: {report_type}")
    
    tpl_path = os.path.join(TEMPLATES_DIR, tpl_file)
    if not os.path.exists(tpl_path):
        raise FileNotFoundError(f"القالب غير موجود: {tpl_path}")
    
    prs = Presentation(tpl_path)
    
    if report_type == "comprehensive":
        fill_comprehensive(prs, state)
    else:
        fill_track(prs, report_type, state)
    
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

if __name__ == "__main__":
    data  = json.loads(sys.stdin.read())
    rtype = data.get("type", "comprehensive")
    state = data.get("state", {})
    result = generate_report(rtype, state)
    sys.stdout.buffer.write(result)
